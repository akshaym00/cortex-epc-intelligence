from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parents[1] / "docs" / "Project_Cortex_Pitch_Deck.pptx"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = "0F172A"; BLUE = "1D4ED8"; CYAN = "0891B2"; TEAL = "0F766E"
GREEN = "16A34A"; ORANGE = "EA580C"; RED = "DC2626"; PURPLE = "7C3AED"
SLATE = "475569"; LIGHT = "F1F5F9"; WHITE = "FFFFFF"; INK = "172033"


def rgb(value):
    return RGBColor.from_string(value)


def shape(slide, x, y, w, h, fill, radius=True, line=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    item = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    item.fill.solid(); item.fill.fore_color.rgb = rgb(fill)
    item.line.color.rgb = rgb(line or fill)
    return item


def text(slide, value, x, y, w, h, size=16, color=INK, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame; frame.clear(); frame.margin_left = 0; frame.margin_right = 0
    frame.margin_top = 0; frame.margin_bottom = 0; frame.vertical_anchor = valign
    for index, line in enumerate(value.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line; paragraph.alignment = align
        paragraph.font.name = "Aptos"; paragraph.font.size = Pt(size)
        paragraph.font.bold = bold; paragraph.font.color.rgb = rgb(color)
    return box


def base(section, heading, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(WHITE)
    shape(slide, 0, 0, 13.333, 0.14, BLUE, False)
    text(slide, section.upper(), 0.55, 0.28, 4.5, 0.25, 9, CYAN, True)
    text(slide, heading, 0.55, 0.72, 12.1, 0.65, 28, INK, True)
    text(slide, subtitle, 0.58, 1.42, 11.9, 0.42, 14, SLATE)
    text(slide, "PROJECT CORTEX", 10.9, 7.12, 1.85, 0.18, 8, "94A3B8", True, PP_ALIGN.RIGHT)
    return slide


# Cover
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb(NAVY)
shape(s, 0, 0, 0.22, 7.5, CYAN, False)
text(s, "PROJECT CORTEX", 0.8, 0.75, 4.8, 0.35, 13, "67E8F9", True)
text(s, "From project documents\nto decision intelligence", 0.8, 1.45, 9.5, 1.6, 38, WHITE, True)
text(s, "An AI reasoning platform for data-centre construction", 0.85, 3.55, 8, 0.45, 18, "CBD5E1")
for i, label in enumerate(["Extract", "Structure", "Reason", "Forecast", "Act"]):
    x = 0.85 + i * 2.25
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(5.05), Inches(.55), Inches(.55))
    circle.fill.solid(); circle.fill.fore_color.rgb = rgb(GREEN if i == 4 else BLUE); circle.line.color.rgb = rgb("64748B")
    text(s, str(i + 1), x, 5.18, .55, .2, 10, WHITE, True, PP_ALIGN.CENTER)
    text(s, label, x - .2, 5.75, .95, .25, 11, WHITE, True, PP_ALIGN.CENTER)
text(s, "ET AI Hackathon 2.0", .85, 6.8, 3.5, .28, 11, "94A3B8")

# Problem
s = base("Problem", "Project data exists. Project intelligence does not.", "Critical dependencies are scattered across RFIs, schedules, reports, submittals and email.")
for i, (metric, label, color) in enumerate([("15K–40K", "schedule line items", BLUE), ("200+", "concurrent contractors", TEAL), ("Hours", "to assess one change manually", ORANGE)]):
    x = .65 + i * 4.15; shape(s, x, 2.15, 3.65, 1.55, LIGHT, True, "CBD5E1")
    text(s, metric, x+.25, 2.42, 3.1, .52, 28, color, True, PP_ALIGN.CENTER)
    text(s, label, x+.25, 3.08, 3.1, .3, 12, SLATE, False, PP_ALIGN.CENTER)
text(s, "A vendor delay is not a document problem.", 1, 4.55, 11.3, .48, 24, INK, True, PP_ALIGN.CENTER)
text(s, "It is a dependency problem: delivery → installation → testing → commissioning → handover", 1, 5.25, 11.3, .4, 16, RED, False, PP_ALIGN.CENTER)

# Solution
s = base("Solution", "A living project model—not another document chatbot", "Cortex transforms unstructured facts into typed, persistent and explainable project intelligence.")
cols = [("AI extraction", "Entities\nRelationships\nIndependent events", BLUE), ("Knowledge graph", "Validated direction\nPersistent structure\nStored dependency paths", CYAN), ("Reasoning", "Impact propagation\nCritical-path forecast\nSpec deviation checks", PURPLE), ("Action", "Quantified recovery\nNamed owners\nExplainable recommendations", GREEN)]
for i, (name, body, color) in enumerate(cols):
    x=.55+i*3.15; shape(s,x,2.1,2.72,3.65,WHITE,True,color)
    shape(s,x+.9,2.4,.92,.72,color,True)
    text(s,str(i+1),x+.9,2.58,.92,.25,15,WHITE,True,PP_ALIGN.CENTER)
    text(s,name,x+.2,3.55,2.32,.36,17,INK,True,PP_ALIGN.CENTER)
    text(s,body,x+.28,4.18,2.16,1.1,12,SLATE,False,PP_ALIGN.CENTER)

# Architecture
s = base("Architecture", "LLMs extract. Deterministic services reason.", "A clear trust boundary makes results defensible under judge questioning.")
nodes=[("Documents",SLATE),("AI extraction",BLUE),("Validation",CYAN),("Living model",TEAL),("Graph",PURPLE),("Reasoning",ORANGE),("Dashboard",GREEN)]
for i,(name,color) in enumerate(nodes):
    x=.35+i*1.84; shape(s,x,2.35,1.48,.92,color,True)
    text(s,name,x+.09,2.62,1.3,.35,11,WHITE,True,PP_ALIGN.CENTER)
for i,(a,b) in enumerate([("Constrained","Enumerated types + known names"),("Validated","Supplier direction + deduplication"),("Calculated","Graph traversal + date arithmetic"),("Explainable","Every impact retains its path")]):
    x=.7+(i%2)*6.15; y=4.15+(i//2)*1.05; shape(s,x,y,5.55,.78,LIGHT,True,"CBD5E1")
    text(s,a,x+.2,y+.18,1.25,.25,12,BLUE,True); text(s,b,x+1.45,y+.18,3.9,.25,12,SLATE)

# Schedule flagship
s = base("Flagship 01", "Predictive schedule risk with defensible date math", "A nine-day vendor delay becomes a quantified critical-path forecast.")
for i,(date,name,color) in enumerate([("15 Jul","Generator delivery",BLUE),("16 Jul","Electrical testing",CYAN),("22 Jul","Commissioning",TEAL)]):
    y=2.12+i*1.22; shape(s,.9,y,.48,.48,color,True); text(s,date,1.65,y,1.2,.3,14,color,True); text(s,name,2.8,y,2.8,.3,14,INK,True)
shape(s,6.5,2,5.6,3.85,"FFF7ED",True,ORANGE)
text(s,"9 DAYS",6.85,2.48,4.9,.65,34,RED,True,PP_ALIGN.CENTER)
text(s,"critical-path exposure",6.85,3.2,4.9,.3,15,SLATE,False,PP_ALIGN.CENTER)
text(s,"Forecast dates shifted deterministically\nSeverity derived from duration\nRecovery scenario: ≈ 4 days",7.05,3.95,4.5,1.15,14,INK,False,PP_ALIGN.CENTER)

# Compliance flagship
s = base("Flagship 02", "Spec compliance with source citations", "A deliberate vendor deviation demonstrates technical assurance beyond document extraction.")
shape(s,.75,2,5.55,3.55,"EFF6FF",True,BLUE); shape(s,7,2,5.55,3.55,"FEF2F2",True,RED)
text(s,"REQUIREMENT",1.1,2.35,2.2,.25,10,BLUE,True); text(s,"UPS efficiency at 50% load",1.1,2.88,4.75,.35,18,INK,True)
text(s,"≥ 96.0%",1.1,3.5,4.75,.62,32,BLUE,True); text(s,"Electrical Specification 26 33 53 §2.4.1 · page 47",1.1,4.6,4.7,.42,11,SLATE)
text(s,"VENDOR SUBMITTAL · DEVIATION",7.35,2.35,3.8,.25,10,RED,True); text(s,"VoltSafe UPS-01",7.35,2.88,4.7,.35,18,INK,True)
text(s,"94.5%",7.35,3.5,2.5,.62,32,RED,True); shape(s,9.72,3.55,1.95,.42,"FEE2E2",True,RED); text(s,"1.5% BELOW SPEC",9.8,3.66,1.78,.18,9,RED,True,PP_ALIGN.CENTER)
text(s,"VoltSafe Submittal VS-442 · page 12",7.35,4.6,4.7,.42,11,SLATE)
text(s,"Action: reject or request a revised submittal before procurement approval.",1,6.15,11.4,.4,16,RED,True,PP_ALIGN.CENTER)

# Scale
s = base("Scalability", "Tested at the brief’s stated line-item scale", "A repeatable benchmark separates measured capability from roadmap claims.")
for i,(metric,label) in enumerate([("15,200","graph entities"),("29,999","relationships"),("14,999","downstream nodes traversed")]):
    x=.6+i*4.22; shape(s,x,2.05,3.78,1.55,NAVY,True); text(s,metric,x+.2,2.35,3.38,.5,29,"67E8F9",True,PP_ALIGN.CENTER); text(s,label,x+.2,3.03,3.38,.25,11,WHITE,False,PP_ALIGN.CENTER)
text(s,"0.096 s\ngraph construction",1.3,4.5,2.2,.85,22,BLUE,True,PP_ALIGN.CENTER)
text(s,"0.013 s\ndependency traversal",5.55,4.5,2.2,.85,22,TEAL,True,PP_ALIGN.CENTER)
text(s,"NEXT\nChunked async extraction\nGraph database\nConcurrent-user load test",9.15,4.42,2.55,1.25,12,ORANGE,True,PP_ALIGN.CENTER)

# Close
s = base("Business value", "From ‘what does the document say?’ to ‘what happens next?’", "Cortex gives project leaders earlier, traceable and actionable visibility.")
for i,(head,body,color) in enumerate([("Earlier","See cascading risk when the notice arrives",BLUE),("Faster","Assess dependencies in seconds",CYAN),("Defensible","Show paths, dates, sources and assumptions",PURPLE),("Actionable","Quantify exposure and recovery options",GREEN)]):
    x=.62+(i%2)*6.2; y=2+(i//2)*1.55; shape(s,x,y,5.65,1.15,LIGHT,True,color); text(s,head,x+.25,y+.22,1.25,.32,18,color,True); text(s,body,x+1.55,y+.24,3.8,.42,13,INK)
shape(s,1.15,5.35,11,.92,NAVY,True); text(s,"Documents → Living model → Knowledge graph → Reasoning → Decisions",1.4,5.66,10.5,.3,18,WHITE,True,PP_ALIGN.CENTER)
text(s,"Project Cortex",4.7,6.62,4,.4,20,BLUE,True,PP_ALIGN.CENTER)

prs.save(OUT)
print(OUT)
